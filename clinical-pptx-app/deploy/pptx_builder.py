"""
在「院內 50 週年模板」上產生簡報。
每一頁都從模板的版面配置（layout）建立 → 一律沿用院內水彩底版與校徽。
文字一律以自訂文字框定位，避免不同 layout 的 placeholder 差異造成重疊。
"""
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

LAYOUT_TITLE = 0      # 封面（左上校徽 + 50 週年標誌）
LAYOUT_CONTENT = 2    # 內頁
MAX_BULLETS_PER_SLIDE = 6
INK = RGBColor(0x33, 0x2E, 0x2A)   # 深暖灰，配合大地色系

SW, SH = 9144000, 5143500          # 16:9 EMU


def _clear_slides(prs):
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        rId = sld.get(qn("r:id"))
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        lst.remove(sld)


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Emu(int(SW * x)), Emu(int(SH * y)),
                                   Emu(int(SW * w)), Emu(int(SH * h)))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _para(tf, text, size, *, bold=False, align=PP_ALIGN.LEFT, first=False, color=INK):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.text = text
    for run in p.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return p


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    tf = _textbox(slide, 0.10, 0.34, 0.80, 0.22)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _para(tf, title or "Journal Reading", 38, bold=True, align=PP_ALIGN.CENTER, first=True)
    if subtitle:
        st = _textbox(slide, 0.12, 0.585, 0.76, 0.16)
        _para(st, subtitle, 17, align=PP_ALIGN.CENTER, first=True)
    return slide


def _add_content_slide(prs, heading, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    ht = _textbox(slide, 0.08, 0.07, 0.84, 0.13)
    _para(ht, heading or "", 28, bold=True, align=PP_ALIGN.CENTER, first=True)
    bt = _textbox(slide, 0.10, 0.24, 0.80, 0.56)
    for i, b in enumerate(bullets):
        _para(bt, "•  " + str(b), 18, align=PP_ALIGN.LEFT, first=(i == 0))
    return slide


def _chunk(items, n=MAX_BULLETS_PER_SLIDE):
    out = [items[i:i + n] for i in range(0, len(items), n)]
    return out or [[]]


def build_deck(deck, template_path="template.pptx"):
    """
    deck = {"title": str, "subtitle": str,
            "slides": [{"title": str, "bullets": [str,...]}, ...]}
    回傳 BytesIO（.pptx）
    """
    prs = Presentation(template_path)
    _clear_slides(prs)
    _add_title_slide(prs, deck.get("title", ""), deck.get("subtitle", ""))
    for s in deck.get("slides", []):
        heading = s.get("title", "")
        bullets = [b for b in (s.get("bullets") or []) if str(b).strip()]
        for gi, g in enumerate(_chunk(bullets)):
            _add_content_slide(prs, heading if gi == 0 else f"{heading}（續）", g)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
